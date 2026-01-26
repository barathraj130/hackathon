# ppt-service/certificate_engine.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Institutional Premium Theme
PRIMARY_COLOR = RGBColor(13, 148, 136) # Teal-600
TEXT_MAIN = RGBColor(30, 41, 59)       # Slate-800
ACCENT_GREY = RGBColor(241, 245, 249)  # Slate-100
WHITE = RGBColor(255, 255, 255)

def create_certificate(name, college, year, dept, role, event_name="HACK@JIT 1.0"):
    prs = Presentation()
    # Use 16:9 Aspect Ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 1. Background Decoration
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.width = 0
    
    # 2. Institutional Border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), prs.slide_width - Inches(0.8), prs.slide_height - Inches(0.8))
    border.fill.background(); border.line.color.rgb = PRIMARY_COLOR; border.line.width = Pt(4)
    
    # Inner thin border
    border2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.55), prs.slide_width - Inches(1.1), prs.slide_height - Inches(1.1))
    border2.fill.background(); border2.line.color.rgb = PRIMARY_COLOR; border2.line.width = Pt(1)

    # 3. Logo
    logo_path = "institution_logo.png"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, (prs.slide_width - Inches(1.5))/2, Inches(0.8), width=Inches(1.5))

    # 4. Content Content
    # "CERTIFICATE OF MERIT/PARTICIPATION"
    tx_cert = slide.shapes.add_textbox(0, Inches(2.6), prs.slide_width, Inches(0.8))
    tf_cert = tx_cert.text_frame; tf_cert.word_wrap = True
    p_cert = tf_cert.paragraphs[0]; p_cert.alignment = PP_ALIGN.CENTER; p_cert.text = "CERTIFICATE OF PARTICIPATION"
    p_cert.font.size = Pt(36); p_cert.font.bold = True; p_cert.font.color.rgb = TEXT_MAIN; p_cert.font.name = 'Arial'

    # "This is to certify that"
    tx_this = slide.shapes.add_textbox(0, Inches(3.4), prs.slide_width, Inches(0.4))
    p_this = tx_this.text_frame.paragraphs[0]; p_this.alignment = PP_ALIGN.CENTER; p_this.text = "This is to certify that"
    p_this.font.size = Pt(16); p_this.font.color.rgb = RGBColor(100, 116, 139); p_this.font.italic = True

    # NAME
    tx_name = slide.shapes.add_textbox(0, Inches(3.9), prs.slide_width, Inches(0.8))
    p_name = tx_name.text_frame.paragraphs[0]; p_name.alignment = PP_ALIGN.CENTER; p_name.text = name.upper()
    p_name.font.size = Pt(44); p_name.font.bold = True; p_name.font.color.rgb = PRIMARY_COLOR

    # Institutional Details
    tx_details = slide.shapes.add_textbox(0, Inches(4.8), prs.slide_width, Inches(0.8))
    tf_d = tx_details.text_frame; tf_d.word_wrap = True
    p_d = tf_d.paragraphs[0]; p_d.alignment = PP_ALIGN.CENTER
    p_d.text = f"of {year} Year, Department of {dept}, {college}"
    p_d.font.size = Pt(16); p_d.font.color.rgb = TEXT_MAIN

    # Participation Text
    tx_event = slide.shapes.add_textbox(0, Inches(5.4), prs.slide_width, Inches(0.8))
    p_e = tx_event.text_frame.paragraphs[0]; p_e.alignment = PP_ALIGN.CENTER
    p_e.text = f"has successfully completed the innovation mission as a {role.capitalize()}\nduring {event_name}, an official institutional hackathon."
    p_e.font.size = Pt(14); p_e.font.color.rgb = RGBColor(71, 85, 105)

    # 5. Bottom Signatures
    # Left: Event Coordinator
    tx_sig1 = slide.shapes.add_textbox(Inches(1.5), Inches(6.4), Inches(3), Inches(0.4))
    p_s1 = tx_sig1.text_frame.paragraphs[0]; p_s1.text = "EVENT COORDINATOR"; p_s1.font.size = Pt(10); p_s1.font.bold = True; p_s1.font.color.rgb = TEXT_MAIN
    line1 = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(6.35), Inches(1.5 + 2.5), Inches(6.35)) # Using rectangle as a line proxy for stability
    
    # Right: Institutional Authority
    tx_sig2 = slide.shapes.add_textbox(prs.slide_width - Inches(4.5), Inches(6.4), Inches(3), Inches(0.4))
    p_s2 = tx_sig2.text_frame.paragraphs[0]; p_s2.alignment = PP_ALIGN.RIGHT; p_s2.text = "PRINCIPAL / DIRECTOR"; p_s2.font.size = Pt(10); p_s2.font.bold = True; p_s2.font.color.rgb = TEXT_MAIN
    
    # Output
    if not os.path.exists('certs_outputs'): os.makedirs('certs_outputs')
    safe_name = name.lower().replace(' ', '_')
    out_path = f"certs_outputs/certificate_{safe_name}.pptx"
    prs.save(out_path)
    return out_path
