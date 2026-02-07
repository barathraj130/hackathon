# ppt-service/generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_pptx(team_name, college, slides_data):
    prs = Presentation()
    
    def add_branding(slide):
        # 1. Top Left - Event Branding
        branding_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(3), Inches(0.4))
        p = branding_box.text_frame.paragraphs[0]
        p.text = "BHARAT BRILLIANT HACKATHON"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(13, 148, 136) # Teal
        
        # 2. Top Right - Logo
        if os.path.exists("hackathon_logo.png"):
            slide.shapes.add_picture("hackathon_logo.png", Inches(8.5), Inches(0.2), width=Inches(1.2))

    def set_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42) # Premium Navy/Black

    # 1. Title Slide (Redesigned with Team Details)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Logo (Centered at top)
    if os.path.exists("hackathon_logo.png"):
        slide.shapes.add_picture("hackathon_logo.png", Inches(4.35), Inches(0.5), height=Inches(0.8))

    # Title Text (Large, Bold, Centered)
    tx_title = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(8.0), Inches(1.0))
    p_title = tx_title.text_frame.paragraphs[0]
    p_title.text = "BRILLIANT BHARAT HACKATHON"
    p_title.font.size = Pt(44); p_title.font.bold = True; p_title.font.name = 'Times New Roman'
    p_title.font.color.rgb = RGBColor(0, 0, 0)
    p_title.alignment = PP_ALIGN.CENTER

    # Subtitle Text (Split into 'Organised by' + College Name)
    tx_subtitle = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(8.0), Inches(1.0))
    tf_s = tx_subtitle.text_frame
    p_org = tf_s.paragraphs[0]
    p_org.text = "Organised by"
    p_org.font.size = Pt(14); p_org.font.bold = False; p_org.font.name = 'Times New Roman'
    p_org.font.color.rgb = RGBColor(0, 0, 0)
    p_org.alignment = PP_ALIGN.CENTER
    
    p_coll = tf_s.add_paragraph()
    p_coll.text = "JANSONS INSTITUTE OF TECHNOLOGY"
    p_coll.font.size = Pt(24); p_coll.font.bold = True; p_coll.font.name = 'Times New Roman'
    p_coll.font.color.rgb = RGBColor(0, 0, 0)
    p_coll.alignment = PP_ALIGN.CENTER

    # Project Name Box
    project_label = slides_data.get('title', {}).get('bullets', ['PROJECT NAME'])[0]
    tx_project = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(8.0), Inches(0.8))
    p_project = tx_project.text_frame.paragraphs[0]
    p_project.text = str(project_label).upper()
    p_project.font.size = Pt(32); p_project.font.bold = True; p_project.font.name = 'Times New Roman'
    p_project.font.color.rgb = RGBColor(13, 148, 136) # Teal
    p_project.alignment = PP_ALIGN.CENTER

    # Team Info (Bottom Centered)
    tx_team = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(7.0), Inches(2.5))
    tf_t = tx_team.text_frame; tf_t.word_wrap = True
    
    p_team = tf_t.paragraphs[0]; p_team.alignment = PP_ALIGN.CENTER
    p_team.text = f"TEAM {team_name.upper()}"; p_team.font.size = Pt(20); p_team.font.bold = True; p_team.font.name = 'Times New Roman'
    
    p_from = tf_t.add_paragraph(); p_from.alignment = PP_ALIGN.CENTER
    team_college = college if college and college != "Institution" else "JANSONS INSTITUTE OF TECHNOLOGY"
    p_from.text = f"from {team_college.upper()}"; p_from.font.size = Pt(14); p_from.font.italic = True; p_from.font.name = 'Times New Roman'
    
    p_leader = tf_t.add_paragraph(); p_leader.alignment = PP_ALIGN.CENTER
    p_leader.text = f"Team Leader: {slides_data.get('leaderName', 'N/A').upper()}"; p_leader.font.size = Pt(16); p_leader.font.bold = True; p_leader.font.name = 'Times New Roman'
    
    p_members = tf_t.add_paragraph(); p_members.alignment = PP_ALIGN.CENTER
    p_members.text = f"MEMBERS: {slides_data.get('memberNames', 'N/A').upper()}"; p_members.font.size = Pt(12); p_members.font.bold = False; p_members.font.name = 'Times New Roman'

    # 2. Add Content Slides
    from pptx.enum.shapes import MSO_SHAPE
    
    for key, data in slides_data.items():
        if key == 'title': continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        add_branding(slide)
        
        # Title - Centered
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
        tf_t = title_box.text_frame
        tf_t.text = data['title']
        p_t = tf_t.paragraphs[0]
        p_t.font.size = Pt(28); p_t.font.bold = True; p_t.font.name = 'Times New Roman'; p_t.alignment = PP_ALIGN.CENTER
        p_t.font.color.rgb = RGBColor(0, 0, 0)
        
        # Content - Box containment
        content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.2))
        content_box.fill.background()
        content_box.line.color.rgb = RGBColor(13, 148, 136) # Institutional Teal
        content_box.line.width = Pt(1.5)
        
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.2)
        
        for point in data['bullets']:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.name = 'Times New Roman'
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(12)

    # Save the file
    file_name = f"ppt_outputs/{team_name.lower().replace(' ', '_')}_pitch_artifact.pptx"
    
    # Ensure directory exists
    if not os.path.exists('ppt_outputs'):
        os.makedirs('ppt_outputs')
        
    prs.save(file_name)
    return file_name