from io import BytesIO
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# Institutional Color Palette
PRIMARY_COLOR = RGBColor(13, 148, 136)  # Teal-600
SECONDARY_COLOR = RGBColor(20, 184, 166)  # Teal-500
TEXT_MAIN = RGBColor(15, 23, 42)  # Slate-900
ACCENT_GREY = RGBColor(241, 245, 249)  # Slate-50
BG_LIGHT = RGBColor(248, 250, 252)  # Slate-100
WHITE = RGBColor(255, 255, 255)
LINE_COLOR = RGBColor(203, 213, 225)  # Slate-300
ERROR_ZONE = RGBColor(239, 68, 68)  # Red-500
SUCCESS_ZONE = RGBColor(34, 197, 94)  # Green-500
WARNING_ZONE = RGBColor(234, 179, 8)  # Amber-500

def disable_shadow(shape):
    """
    POLARIS PROTOCOL: Hard-disable all shadows for clean, modern aesthetics.
    """
    try:
        if hasattr(shape, 'shadow'):
            shape.shadow.inherit = False
            # Some versions of pptx require setting visibility
            try: shape.shadow.visible = False
            except: pass
    except:
        pass
    return shape

def set_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_header(slide, title="SLIDE TITLE"):
    # VIKSIT BHARAT GLOBAL LOGO (Top Left)
    if os.path.exists("viksit_bharat_global.png"):
        vb_logo = slide.shapes.add_picture("viksit_bharat_global.png", Inches(0.2), Inches(0.2), height=Inches(0.6))
        disable_shadow(vb_logo)

    # SLIDE TITLE (Shifted right for logo)
    header_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(7.0), Inches(0.5))
    disable_shadow(header_box)
    p = header_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(18); p.font.bold = True; p.font.name = 'Arial Black'
    p.font.color.rgb = TEXT_MAIN
    p.alignment = PP_ALIGN.LEFT

    # TOP RIGHT LOGO (Institutional Branding)
    if os.path.exists("institution_logo.png"):
        logo = slide.shapes.add_picture("institution_logo.png", Inches(8.8), Inches(0.2), height=Inches(0.6))
        disable_shadow(logo)

def add_clean_box(slide, text, x, y, w, h, sz, bold=False, txt_color=TEXT_MAIN, border_color=None, bg_color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    disable_shadow(box)
    if bg_color:
        box.fill.solid(); box.fill.fore_color.rgb = bg_color
    else:
        box.fill.background()
    if border_color:
        box.line.color.rgb = border_color; box.line.width = Pt(1)
    else:
        box.line.width = 0
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = txt_color


def create_expert_deck(team_name, college, data):
    prs = Presentation()
    
    # ... (Cover Slide Logic - Unchanged)
    
    # 1. PREMIUM COVER SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    
    # Large centered Logo on cover
    if os.path.exists("institution_logo.png"):
        slide.shapes.add_picture("institution_logo.png", Inches(4.25), Inches(0.5), height=Inches(1.2))

    # Viksit Bharat Branding - Top Left
    if os.path.exists("viksit_bharat_global.png"):
        vb_img = slide.shapes.add_picture("viksit_bharat_global.png", Inches(0.5), Inches(0.4), height=Inches(0.8))
        disable_shadow(vb_img)
        # Optional text label for Viksit Bharat
        tx_vb = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(2), Inches(0.3))
        p_vb = tx_vb.text_frame.paragraphs[0]
        p_vb.text = "VIKSIT BHARAT @ 2047"; p_vb.font.size = Pt(8); p_vb.font.bold = True; p_vb.font.color.rgb = PRIMARY_COLOR
        p_vb.alignment = PP_ALIGN.LEFT

    # High-impact Project Title - REDUCED FONT TO PREVENT OVERLAP
    tx_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.8))
    p_title = tx_title.text_frame.paragraphs[0]
    p_title.text = data.get('projectName', 'VENTURE PROTOTYPE').upper()
    p_title.font.size = Pt(38); p_title.font.bold = True; p_title.font.color.rgb = TEXT_MAIN; p_title.alignment = PP_ALIGN.CENTER
    p_title.line_spacing = 1.0
    
    # Elegant Underline on cover - ADJUSTED POSITION
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.5), Inches(3.7), Inches(6.5), Inches(3.7))
    line.line.color.rgb = TEXT_MAIN; line.line.width = Pt(2)

    # Detailed Personnel Info - LOWERED TO PREVENT OVERLAP
    tx_id = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2.5))
    tf_d = tx_id.text_frame; tf_d.word_wrap = True
    
    p1 = tf_d.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    p1.text = f"TEAM {team_name.upper()}"; p1.font.size = Pt(24); p1.font.bold = True; p1.font.color.rgb = TEXT_MAIN
    p1.space_after = Pt(8)
    
    p2 = tf_d.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    p2.text = f"from {college.upper()}"; p2.font.size = Pt(14); p2.font.italic = True; p2.font.color.rgb = TEXT_MAIN
    p2.space_after = Pt(12)
    
    p3 = tf_d.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    p3.text = f"Team Leader: {data.get('leaderName', 'N/A').upper()}"; p3.font.size = Pt(18); p3.font.bold = True; p3.font.color.rgb = TEXT_MAIN; p3.font.name = 'Poppins'
    p3.space_after = Pt(8)
    
    p4 = tf_d.add_paragraph(); p4.alignment = PP_ALIGN.CENTER
    p4.text = f"MEMBERS: {data.get('memberNames', 'N/A').upper()}"; p4.font.size = Pt(11); p4.font.bold = False; p4.font.color.rgb = TEXT_MAIN; p4.font.name = 'Inter'
    
    add_footer(slide)

    modules = [
        ("Background", lambda s: draw_strategic(s, data)),
        ("The Problem", lambda s: draw_problem(s, data)),
        ("Impact", lambda s: draw_impact(s, data)),
        ("People Involved", lambda s: draw_stakeholders(s, data)),
        ("Our User", lambda s: draw_persona(s, data)),
        ("What's Missing?", lambda s: draw_gap(s, data)),
        ("Our Solution", lambda s: draw_solution_statement(s, data)),
        ("How it Works", lambda s: draw_solution_flow(s, data)),
        ("The Plan", lambda s: draw_lean(s, data)),
        ("Growth", lambda s: draw_balloon(s, data)),
        ("Competition", lambda s: draw_market_matrix(s, data)),
        ("Market Size", lambda s: draw_market_sizing(s, data)),
        ("Business Model", lambda s: draw_revenue(s, data)),
        ("Budget", lambda s: draw_fiscal(s, data)),
        ("The Future", lambda s: draw_vision(s, data))
    ]

    for title, fn in modules:
        s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s)
        add_header(s, title); fn(s)

    # CLOSURE (Unchanged)
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide)
    if os.path.exists("institution_logo.png"):
        slide.shapes.add_picture("institution_logo.png", Inches(4.25), Inches(0.5), height=Inches(1.2))
    
    # Viksit Bharat branding on closure
    if os.path.exists("viksit_bharat_global.png"):
        vb_close = slide.shapes.add_picture("viksit_bharat_global.png", Inches(0.5), Inches(0.4), height=Inches(0.8))
        disable_shadow(vb_close)

    tx = slide.shapes.add_textbox(Inches(0), Inches(3.2), Inches(10), Inches(1.5))
    p = tx.text_frame.paragraphs[0]; p.text = "THANK YOU."; p.font.size = Pt(64); p.font.bold = True; p.font.color.rgb = TEXT_MAIN; p.alignment = PP_ALIGN.CENTER
    add_footer(slide)

    if not os.path.exists('ppt_outputs'): os.makedirs('ppt_outputs')
    out = f"ppt_outputs/{team_name.lower().replace(' ', '_')}_pitch_artifact.pptx"
    prs.save(out); return out

# --- DRAWERS ---

def draw_strategic(slide, data):
    add_clean_box(slide, "DOMAIN", Inches(0.5), Inches(1.8), Inches(9), Inches(0.35), 11, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s2_domain', 'N/A'), Inches(0.5), Inches(2.2), Inches(9), Inches(0.8), 14)
    add_clean_box(slide, "OPERATIONAL CONTEXT", Inches(0.5), Inches(3.2), Inches(9), Inches(0.35), 11, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s2_context', 'N/A'), Inches(0.5), Inches(3.6), Inches(9), Inches(1.8), 12)
    add_clean_box(slide, "ROOT CATALYST", Inches(0.5), Inches(5.6), Inches(9), Inches(0.35), 11, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s2_rootReason', 'N/A'), Inches(0.5), Inches(6.0), Inches(9), Inches(0.9), 12)

def draw_problem(slide, data):
    add_clean_box(slide, "CORE PROBLEM", Inches(0.5), Inches(1.8), Inches(9), Inches(0.4), 12, True, ERROR_ZONE, None, BG_LIGHT)
    add_clean_box(slide, data.get('s3_coreProblem', 'N/A'), Inches(0.5), Inches(2.3), Inches(9), Inches(2.0), 16, False, TEXT_MAIN, LINE_COLOR, WHITE)
    add_clean_box(slide, "AFFECTED PERSONNEL", Inches(0.5), Inches(4.5), Inches(4.4), Inches(0.35), 11, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s3_affected', 'N/A'), Inches(0.5), Inches(4.9), Inches(4.4), Inches(1.8), 12)
    add_clean_box(slide, "CRITICAL GRAVITY", Inches(5.1), Inches(4.5), Inches(4.4), Inches(0.35), 11, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s3_whyItMatters', 'N/A'), Inches(5.1), Inches(4.9), Inches(4.4), Inches(1.8), 12)

def draw_stakeholders(slide, data):
    for i, (label, key) in enumerate([("PRIMARY SEGMENT", 's5_primaryUsers'), ("SECONDARY SEGMENT", 's5_secondaryUsers')]):
        y_pos = 1.8 + i * 2.6
        add_clean_box(slide, label, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.35), 12, True, PRIMARY_COLOR, BG_LIGHT, BG_LIGHT)
        add_clean_box(slide, data.get(key, 'N/A'), Inches(0.5), Inches(y_pos + 0.4), Inches(9), Inches(1.8), 14)

def draw_persona(slide, data):
    coords = [(0.5, 1.8), (5.1, 1.8), (0.5, 4.4), (5.1, 4.4)]
    titles = ["PERSONAL INFO", "CHALLENGES", "PROFESSIONAL GOALS", "SUCCESS FACTORS"]
    vals = [
        f"Name: {data.get('s6_customerName','X')}\nAge: {data.get('s6_customerAge', 'X')}\nLoc: {data.get('s6_customerLocation','X')}",
        data.get('s6_pains', 'N/A'),
        data.get('s6_goals', 'N/A'),
        data.get('s6_howWeHelp', 'N/A')
    ]
    for i, (x, y) in enumerate(coords):
        add_clean_box(slide, titles[i], Inches(x), Inches(y), Inches(4.4), Inches(0.35), 11, True, PRIMARY_COLOR, None, BG_LIGHT)
        add_clean_box(slide, vals[i], Inches(x), Inches(y + 0.4), Inches(4.4), Inches(2.0), 11)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.4), Inches(3.7), Inches(1.2), Inches(1.2))
    c.fill.solid(); c.fill.fore_color.rgb = PRIMARY_COLOR; c.line.color.rgb = WHITE; c.line.width = Pt(2)
    add_text_box_centered(slide, data.get('s6_customerName', 'PERSONA').upper()[:10], 4.4, 4.15, 1.2, 0.3, 9, True, WHITE)


def draw_impact(slide, data):
    # Professional Visualization Grid
    x0, y0, w, h = 1.0, 1.8, 5.0, 4.8
    # Axes
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0+h), Inches(x0+w+0.2), Inches(y0+h)).line.color.rgb = TEXT_MAIN
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0+h), Inches(x0), Inches(y0-0.2)).line.color.rgb = TEXT_MAIN
    
    # Quadrant Lines (Target Centers)
    q1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0+w/2), Inches(y0), Inches(x0+w/2), Inches(y0+h))
    q1.line.color.rgb = LINE_COLOR; q1.line.width = Pt(1); q1.line.dash_style = 2
    q2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0+h/2), Inches(x0+w), Inches(y0+h/2))
    q2.line.color.rgb = LINE_COLOR; q2.line.width = Pt(1); q2.line.dash_style = 2
    
    # Axis Labels (Watermarks removed as requested)
    add_text_box_simple(slide, "CRITICAL IMPACT ⭡", x0+0.1, y0 - 0.4, 2.0, 0.4, 10, True, ERROR_ZONE)
    add_text_box_simple(slide, "FREQUENCY ⭢", x0 + w - 1.0, y0+h+0.1, 1.5, 0.4, 10, True, PRIMARY_COLOR)

    m = {"Low": 1, "Medium": 2, "High": 3, "Rare": 1, "Occasional": 2, "Frequent": 3}
    pts = [p for p in data.get('s4_painPoints', []) if p.get('point')]
    
    for i, p in enumerate(pts[:8]):
        ix = m.get(p.get('freq'), 2); iy = m.get(p.get('impact'), 2)
        px = x0 + (ix/3.8)*w; py = (y0+h) - (iy/3.8)*h
        
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px-0.18), Inches(py-0.18), Inches(0.36), Inches(0.36))
        disable_shadow(dot)
        
        # COLOR-CODED IMPACT LOGIC
        imp_val = p.get('impact', 'Medium')
        dot_color = ERROR_ZONE # Default High
        if imp_val == 'Low': dot_color = SUCCESS_ZONE
        elif imp_val == 'Medium': dot_color = WARNING_ZONE
        
        dot.fill.solid(); dot.fill.fore_color.rgb = dot_color; dot.line.color.rgb = WHITE; dot.line.width = Pt(1.5)
        
        tx_dot = slide.shapes.add_textbox(Inches(px-0.18), Inches(py-0.18), Inches(0.36), Inches(0.36))
        disable_shadow(tx_dot)
        p_dot = tx_dot.text_frame.paragraphs[0]; p_dot.text = str(i+1); p_dot.font.size=Pt(10); p_dot.font.bold=True; p_dot.font.color.rgb=WHITE; p_dot.alignment=PP_ALIGN.CENTER
        
        ly = y0 + (i * 0.6)
        leg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(ly), Inches(3.2), Inches(0.5))
        disable_shadow(leg)
        leg.fill.solid(); leg.fill.fore_color.rgb = dot_color # Match legend to dot
        leg.line.color.rgb = WHITE; leg.line.width = Pt(0.5)
        p_leg = leg.text_frame.paragraphs[0]; p_leg.text = f"{i+1}. {p['point'][:60]}"; p_leg.font.size=Pt(9); p_leg.font.bold=True; p_leg.font.color.rgb=WHITE

def draw_prototype(slide, data):
    # Fetch images from data s8_5_img1, s8_5_img2, s8_5_img3
    images = [data.get(f's8_5_img{i}') for i in range(1, 4)]
    images = [img for img in images if img and img.startswith('http')]
    
    if not images:
        add_text_box_centered(slide, "PROTOTYPE EVIDENCE NODES NOT UPLOADED", 1.0, 3.0, 8.0, 1.0, 24, True, LINE_COLOR)
        return

    # Dynamic layout based on count
    import io
    count = len(images)
    
    # Layout Configs: (left, top, width, height)
    layouts = {
        1: [(1.0, 1.8, 8.0, 4.5)],
        2: [(0.5, 2.2, 4.4, 3.3), (5.1, 2.2, 4.4, 3.3)],
        3: [(0.5, 2.2, 2.8, 3.3), (3.6, 2.2, 2.8, 3.3), (6.7, 2.2, 2.8, 3.3)]
    }
    
    cfgs = layouts.get(count, layouts[3])
    
    for i, url in enumerate(images[:3]):
        l, t, w, h = cfgs[i]
        try:
            # Download image
            response = requests.get(url, timeout=5)
            if response.status_code == 200:
                img_stream = io.BytesIO(response.content)
                pic = slide.shapes.add_picture(img_stream, Inches(l), Inches(t), width=Inches(w))
                disable_shadow(pic)
                # Add border
                pic.line.color.rgb = PRIMARY_COLOR
                pic.line.width = Pt(2)
                # Add Label
                add_text_box_centered(slide, f"EVIDENCE NODE {chr(65+i)}", l, t+h+0.1, w, 0.3, 10, True, SECONDARY_COLOR)
            else:
                add_clean_box(slide, f"IMG LOAD FAIL: {url[:30]}...", Inches(l), Inches(t), Inches(w), Inches(h), 10, False, ERROR_ZONE)
        except Exception as e:
            print(f"Failed to load image {url}: {e}")
            add_clean_box(slide, "IMAGE STREAM INTERRUPTED", Inches(l), Inches(t), Inches(w), Inches(h), 12, True, ERROR_ZONE, BG_LIGHT, BG_LIGHT)

# ... (Rest of drawers: draw_solution_flow, draw_lean, etc. Unchanged)


# --- HELPER FUNCTIONS ---

def add_text_box_simple(slide, text, x, y, w, h, sz, b=False, cl=TEXT_MAIN):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    disable_shadow(tx)
    p = tx.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = cl

def add_text_box_centered(slide, text, x, y, w, h, sz, b, cl):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    disable_shadow(tx)
    p = tx.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = cl; p.alignment = PP_ALIGN.CENTER

def add_footer(slide, text=""):
    # Footer disabled as per watermark removal request
    pass

def draw_gap(slide, data):
    pts = [((0.5, 1.8), "STATUS QUO", 's7_alternatives'), ((5.1, 1.8), "SYSTEMIC GAPS", 's7_limitations'), ((0.5, 4.4), "VALUE GAINS", 's7_gainCreators'), ((5.1, 4.4), "PAIN RELIEF", 's7_painKillers')]
    for (x,y), t, k in pts:
        add_clean_box(slide, t, Inches(x), Inches(y), Inches(4.4), Inches(0.35), 11, True, PRIMARY_COLOR, None, BG_LIGHT)
        add_clean_box(slide, data.get(k, 'N/A'), Inches(x), Inches(y+0.4), Inches(4.4), Inches(2.0), 11)

def draw_solution_statement(slide, data):
    add_clean_box(slide, "THE VENTURE UNVEILED", Inches(0.5), Inches(1.6), Inches(9), Inches(0.4), 12, True, PRIMARY_COLOR, BG_LIGHT, BG_LIGHT)
    h_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1), Inches(9), Inches(2.8))
    h_box.fill.solid(); h_box.fill.fore_color.rgb = WHITE; h_box.line.color.rgb = PRIMARY_COLOR; h_box.line.width = Pt(2)
    p = h_box.text_frame.paragraphs[0]; p.text = data.get('s8_solution', 'N/A'); p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = TEXT_MAIN; p.alignment = PP_ALIGN.CENTER
    add_clean_box(slide, "CORE TECHNOLOGY ARCHITECTURE", Inches(0.5), Inches(5.1), Inches(9), Inches(0.4), 12, True, SECONDARY_COLOR, BG_LIGHT, BG_LIGHT)
    add_clean_box(slide, data.get('s8_coreTech', 'N/A'), Inches(0.5), Inches(5.6), Inches(9), Inches(1.1), 16, False, PRIMARY_COLOR)

def draw_solution_flow(slide, data):
    add_clean_box(slide, data.get('s9_oneline', 'N/A'), Inches(0.5), Inches(1.5), Inches(9), Inches(0.6), 22, True, PRIMARY_COLOR, BG_LIGHT, BG_LIGHT)
    sps = [s for s in data.get('s9_flowSteps', []) if s.strip()][:6]
    for i, s in enumerate(sps):
        x, y = 0.5 + (i%3)*3.2, 3.2 + (i//3)*1.6
        add_clean_box(slide, f"{i+1}. {s}", Inches(x), Inches(y), Inches(2.8), Inches(1.3), 10, True, TEXT_MAIN, PRIMARY_COLOR, ACCENT_GREY)

def draw_lean(slide, data):
    w = 1.9; pillars = [('PROBLEM','s10_leanProblem',0.4,4.0), ('SOLUTION','s10_leanSolution',0.4+w,2.0), ('USP','s10_leanUSP',0.4+2*w,4.0), ('ADVANAGE','s10_leanUnfair',0.4+3*w,2.0), ('SEGMENTS','s10_leanSegments',0.4+4*w,4.0)]
    for t,k,x,h in pillars:
        add_clean_box(slide, t, Inches(x), Inches(1.6), Inches(w-0.1), Inches(0.35), 9, True, PRIMARY_COLOR, None, BG_LIGHT)
        add_clean_box(slide, data.get(k,'N/A'), Inches(x), Inches(2.0), Inches(w-0.1), Inches(h-0.4), 8)
    for t,k,x,y in [("METRICS","s10_leanMetrics",0.4+w,3.6), ("CHANNELS","s10_leanChannels",0.4+3*w,3.6)]:
        add_clean_box(slide, t, Inches(x), Inches(y), Inches(w-0.1), Inches(0.35), 8, True, PRIMARY_COLOR, None, BG_LIGHT)
        add_clean_box(slide, data.get(k,'N/A'), Inches(x), Inches(y+0.4), Inches(w-0.1), Inches(1.6), 8)
    for t,k,x in [("COSTS","s10_leanCosts",0.4), ("REVENUE","s10_leanRevenue",0.4+w*3)]:
        add_clean_box(slide, t, Inches(x), Inches(5.7), Inches(w*2), Inches(0.35), 9, True, PRIMARY_COLOR, ACCENT_GREY, ACCENT_GREY)
        add_clean_box(slide, data.get(k,'N/A'), Inches(x), Inches(6.1), Inches(w*2), Inches(0.9), 8, False, TEXT_MAIN, ACCENT_GREY, ACCENT_GREY)

def draw_balloon(slide, data):
    env = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(1.5), Inches(3.6), Inches(3.6))
    env.fill.solid(); env.fill.fore_color.rgb = PRIMARY_COLOR; env.line.color.rgb = SECONDARY_COLOR; env.line.width = Pt(1)
    add_text_box_centered(slide, "LIFTS (DRIVERS)", 3.4, 2.0, 3.2, 0.4, 12, True, WHITE)
    ls = "\n".join([f"• {x}" for x in data.get('s11_lifts', []) if x.strip()][:4])
    add_text_box_centered(slide, ls, 3.4, 2.4, 3.2, 1.5, 10, False, WHITE)
    bk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(5.6), Inches(1.6), Inches(1.0))
    bk.fill.solid(); bk.fill.fore_color.rgb = SECONDARY_COLOR; bk.line.width = 0
    add_text_box_centered(slide, "VENTURE CORE", 4.2, 5.8, 1.6, 0.4, 11, True, WHITE)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.7), Inches(4.8), Inches(4.2), Inches(5.6)).line.color.rgb=SECONDARY_COLOR
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.3), Inches(4.8), Inches(5.8), Inches(5.6)).line.color.rgb=SECONDARY_COLOR
    add_clean_box(slide, "PULLS (ANCHORS)", Inches(0.4), Inches(4.0), Inches(2.7), Inches(0.35), 11, True, ERROR_ZONE, BG_LIGHT, BG_LIGHT)
    pl = "\n".join([f"• {x}" for x in data.get('s11_pulls', []) if x.strip()][:4])
    add_clean_box(slide, pl, Inches(0.4), Inches(4.4), Inches(2.7), Inches(1.8), 10)
    add_clean_box(slide, "OUTCOMES (ALTITUDE)", Inches(6.9), Inches(4.0), Inches(2.7), Inches(0.35), 11, True, PRIMARY_COLOR, BG_LIGHT, BG_LIGHT)
    os = "\n".join([f"• {x}" for x in data.get('s11_outcomes', []) if x.strip()][:4])
    add_clean_box(slide, os, Inches(6.9), Inches(4.4), Inches(2.7), Inches(1.8), 10)

def draw_market_matrix(slide, data):
    rows = 4; cols = 4
    t = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(9), Inches(5)).table
    hdrs = ["FEATURE / METRIC", "COMPETITOR 1", "COMPETITOR 2", "OUR VENTURE"]
    for i, h in enumerate(hdrs):
        c = t.cell(0, i); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = PRIMARY_COLOR
        p = c.text_frame.paragraphs[0]; p.font.size=Pt(11); p.font.bold=True; p.font.color.rgb=WHITE
    
    f_rows = ["Market Depth", "Pricing Model", "Feature Richness", "Future Readiness"]
    comps = data.get('s12_competitors', [])
    our = data.get('s12_ourVenture', {})
    
    for r in range(1, rows):
        t.cell(r, 0).text = f_rows[r-1]
        
        # Competitor 1 data
        if len(comps) > 0:
            if r == 1:  # Market Depth
                t.cell(r, 1).text = comps[0].get('strength', 'N/A')
            elif r == 2:  # Pricing Model
                t.cell(r, 1).text = comps[0].get('pricingModel', 'Baseline')
            elif r == 3:  # Feature Richness
                t.cell(r, 1).text = comps[0].get('featureRichness', 'Baseline')
            else:  # Future Readiness
                t.cell(r, 1).text = "Baseline"
        
        # Competitor 2 data
        if len(comps) > 1:
            if r == 1:  # Market Depth
                t.cell(r, 2).text = comps[1].get('strength', 'N/A')
            elif r == 2:  # Pricing Model
                t.cell(r, 2).text = comps[1].get('pricingModel', 'Baseline')
            elif r == 3:  # Feature Richness
                t.cell(r, 2).text = comps[1].get('featureRichness', 'Baseline')
            else:  # Future Readiness
                t.cell(r, 2).text = "Baseline"
        
        # Our Venture data
        if r == 1:  # Market Depth
            t.cell(r, 3).text = our.get('strength', 'N/A')
        elif r == 2:  # Pricing Model
            t.cell(r, 3).text = our.get('pricingModel', 'Disruptive')
        elif r == 3:  # Feature Richness
            t.cell(r, 3).text = our.get('featureRichness', 'Disruptive')
        else:  # Future Readiness
            t.cell(r, 3).text = "Disruptive"
        
        # Apply styling
        for c in range(cols):
            p = t.cell(r, c).text_frame.paragraphs[0]; p.font.size = Pt(10)
            if r % 2 == 0:
                t.cell(r, c).fill.solid(); t.cell(r, c).fill.fore_color.rgb = ACCENT_GREY

def draw_market_sizing(slide, data):
    # Professional Bullseye Logic (TAM > SAM > SOM)
    cx, cy = 3.8, 3.8 # Center Point
    configs = [ (4.4, PRIMARY_COLOR, 0.8), (3.2, SECONDARY_COLOR, 0.6), (2.0, ERROR_ZONE, 0.4) ]
    for d, cl, tr in configs:
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-d/2), Inches(cy-d/2), Inches(d), Inches(d))
        disable_shadow(sh)
        sh.fill.solid(); sh.fill.fore_color.rgb = cl; sh.fill.transparency = tr
        sh.line.color.rgb = cl; sh.line.width = Pt(1.5)
    
    lbls = [("TAM", 's13_tam', PRIMARY_COLOR, 2.0), ("SAM", 's13_sam', SECONDARY_COLOR, 2.8), ("SOM", 's13_som', ERROR_ZONE, 3.6)]
    for i, (name, key, col, y_pos) in enumerate(lbls):
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx+0.5), Inches(y_pos), Inches(7.0), Inches(y_pos)).line.color.rgb = col
        add_text_box_simple(slide, f"{name}: {data.get(key, 'N/A')}", 7.1, y_pos-0.2, 2.5, 0.4, 15, True, col)

    add_clean_box(slide, "VALUATION LOGIC & SOURCE DATA", Inches(0.5), Inches(6.1), Inches(9), Inches(0.35), 10, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s13_marketLogic', 'N/A'), Inches(0.5), Inches(6.45), Inches(9), Inches(0.65), 10)

def draw_revenue(slide, data):
    m = [("PRIMARY STREAM", 's14_primaryStream', 0.5, 1.8), ("SECONDARY STREAM", 's14_secondaryStream', 5.1, 1.8), ("PRICING LOGIC", 's14_pricingStrategy', 0.5, 4.4), ("ECONOMIC LOGIC", 's14_revenueLogic', 5.1, 4.4)]
    for lb, k, x, y in m:
        add_clean_box(slide, lb, Inches(x), Inches(y), Inches(4.4), Inches(0.35), 12, True, PRIMARY_COLOR, None, BG_LIGHT)
        add_clean_box(slide, data.get(k, 'N/A'), Inches(x), Inches(y+0.4), Inches(4.4), Inches(2.0), 11)

def draw_fiscal(slide, data):
    als = [a for a in data.get('s15_allocations', []) if a.get('category')]
    if not als: als = [{"category": "SYSTEMIC DEVELOPMENT", "amount": "0"}]
    
    total = 0
    for a in als:
        try:
            val = str(a.get('amount', '0')).replace(',', '').replace('₹', '').strip()
            total += float(val) if val else 0
        except: pass

    rows = len(als) + 2 # Header + Data + Total
    t = slide.shapes.add_table(rows, 2, Inches(1), Inches(2.0), Inches(8), Inches(min(5, rows*0.6))).table
    t.cell(0,0).text = "ALLOCATION NODE"; t.cell(0,1).text = "VALUATION / PURPOSE (₹)"
    
    for r in range(rows):
        for c in range(2):
            cell = t.cell(r,c); cell.fill.solid()
            # Styling branching
            is_header = (r == 0)
            is_total = (r == rows - 1)
            
            if is_header:
                cell.fill.fore_color.rgb = PRIMARY_COLOR
            elif is_total:
                cell.fill.fore_color.rgb = SECONDARY_COLOR
            else:
                cell.fill.fore_color.rgb = ACCENT_GREY if r % 2 == 0 else WHITE
                
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14 if is_total else 11)
            p.font.bold = is_header or is_total
            p.font.color.rgb = WHITE if (is_header or is_total) else TEXT_MAIN
            
            if r > 0 and r < rows - 1:
                idx = r - 1
                cell.text = als[idx]['category'].upper() if c == 0 else f"₹ {als[idx]['amount']}"
            elif is_total:
                cell.text = "TOTAL FISCAL ALLOCATION" if c == 0 else f"₹ {total:,.2f}"

def draw_vision(slide, data):
    add_clean_box(slide, "MACRO IMPACT", Inches(0.5), Inches(1.8), Inches(9), Inches(0.35), 12, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s16_socialEconomic', 'N/A'), Inches(0.5), Inches(2.2), Inches(9), Inches(2.3), 12)
    add_clean_box(slide, "FUTURE TRAJECTORY", Inches(0.5), Inches(4.8), Inches(9), Inches(0.35), 12, True, PRIMARY_COLOR, None, BG_LIGHT)
    add_clean_box(slide, data.get('s16_vision', 'N/A'), Inches(0.5), Inches(5.2), Inches(9), Inches(1.5), 12)
